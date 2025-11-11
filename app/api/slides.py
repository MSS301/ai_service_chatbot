from fastapi import APIRouter, HTTPException, UploadFile, File, Form, Response
from typing import Any, Dict, List, Optional
from app.models.rag_model import (
    SlidesGPTRequest, SlidesGPTResponse,
    TemplateSlidesRequest, TemplateSlidesResponse
)
from app.core.config import SLIDES_BASE_URL, SLIDESGPT_API_KEY, OPENAI_API_KEY
import requests, uuid, os
from pydantic import BaseModel
from app.repositories.content_repository import ContentRepository
from app.repositories.template_repository import SlideTemplateRepository
from openai import OpenAI
import re
import yaml
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

router = APIRouter()

@router.post("/slidesgpt", response_model=SlidesGPTResponse)
def create_with_slidesgpt(req: SlidesGPTRequest):
    """
    Tạo slide qua SlidesGPT (proxy).
    Body: { "prompt": "..." }
    Trả về: { id, embed, download }
    """
    base = SLIDES_BASE_URL.rstrip("/")
    api_key = SLIDESGPT_API_KEY or os.getenv("SLIDESGPT_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="SLIDESGPT_API_KEY not configured")

    url = f"{base}/v1/presentations/generate"
    try:
        r = requests.post(
            url,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={"prompt": req.prompt},
            timeout=120,
        )
        if r.status_code >= 400:
            raise HTTPException(status_code=r.status_code, detail=f"SlidesGPT error: {r.text}")
        data = r.json()
        return {
            "id": data.get("id", uuid.uuid4().hex),
            "embed": data.get("embed"),
            "download": data.get("download"),
        }
    except requests.RequestException as e:
        raise HTTPException(status_code=502, detail=f"SlidesGPT request failed: {e}")

class SlidesGPTFromContentRequest(BaseModel):
    content_id: str
    created_by: str | None = None

@router.post("/gpt", response_model=SlidesGPTResponse)
def create_with_slidesgpt_from_content(req: SlidesGPTFromContentRequest):
    """
    Tạo slide qua SlidesGPT bằng content_id.
    Hệ thống tự lấy content_text đã sinh làm prompt.
    """
    # Load content
    crepo = ContentRepository()
    doc = crepo.get_by_id(req.content_id)
    if not doc:
        raise HTTPException(status_code=404, detail="content_id not found")
    prompt = doc.get("content_text", "")
    if not prompt:
        raise HTTPException(status_code=400, detail="content_text is empty for this content_id")

    # Call SlidesGPT
    base = SLIDES_BASE_URL.rstrip("/")
    api_key = SLIDESGPT_API_KEY or os.getenv("SLIDESGPT_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="SLIDESGPT_API_KEY not configured")

    url = f"{base}/v1/presentations/generate"
    try:
        r = requests.post(
            url,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={"prompt": prompt},
            timeout=120,
        )
        if r.status_code >= 400:
            raise HTTPException(status_code=r.status_code, detail=f"SlidesGPT error: {r.text}")
        data = r.json()
        resp = {
            "id": data.get("id", uuid.uuid4().hex),
            "embed": data.get("embed"),
            "download": data.get("download"),
        }
        # Save to DB under this content_id
        crepo.save_slidesgpt(req.content_id, resp, created_by=req.created_by)
        return resp
    except requests.RequestException as e:
        raise HTTPException(status_code=502, detail=f"SlidesGPT request failed: {e}")


@router.post("/template", response_model=TemplateSlidesResponse)
def create_with_template(req: TemplateSlidesRequest):
    """
    Tạo slide theo khung template (JSON) để client render PPT/HTML.
    Body: { title, outline, theme? }
    """
    outline = req.outline or {}
    sections = outline.get("sections", [])
    slides = []
    slides.append({"type": "title", "title": req.title, "subtitle": req.theme or ""})
    for sec in sections:
        slides.append({
            "type": "content",
            "title": sec.get("title", ""),
            "bullets": sec.get("bullets", []),
            "examples": sec.get("examples", []),
        })
    slides.append({"type": "closing", "title": "Tổng kết", "bullets": ["Câu hỏi?", "Bài tập/vận dụng"]})
    return {"slides": slides}


class TemplateYAMLFromContentRequest(BaseModel):
    content_id: str
    created_by: str | None = None


class TemplateYAMLResponse(BaseModel):
    content_yaml_id: str
    yaml: str


@router.post("/template/yaml", response_model=TemplateYAMLResponse)
def generate_template_yaml_from_content(req: TemplateYAMLFromContentRequest):
    """
    Sinh YAML slide từ content_id theo schema:
    slides: [ {layout, title, bullets: [...] } ]
    meta: { deck_title, author }
    """
    crepo = ContentRepository()
    doc = crepo.get_by_id(req.content_id)
    if not doc:
        raise HTTPException(status_code=404, detail="content_id not found")
    content_text = doc.get("content_text", "")
    if not content_text:
        raise HTTPException(status_code=400, detail="content_text is empty for this content_id")
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY not configured")

    client = OpenAI(api_key=OPENAI_API_KEY)
    system_msg = (
        "Chỉ trả về YAML hợp lệ, KHÔNG kèm giải thích, KHÔNG code fence. "
        "Yêu cầu 5–10 slide. Mỗi phần tử có các khóa: slide (số thứ tự), title (ngắn gọn), content (đoạn nhiều dòng)."
    )
    user_msg = (
        "Chuyển nội dung sau thành YAML tạo slide theo schema:\n\n"
        "slides:\n"
        "  - slide: 1\n"
        "    title: \"Tiêu đề ngắn\"\n"
        "    content: |\n"
        "      Dòng 1\n"
        "      Dòng 2\n"
        "  - slide: 2\n"
        "    title: \"Tiêu đề khác\"\n"
        "    content: |\n"
        "      Ý chính 1\n"
        "      Ý chính 2\n"
        "meta:\n"
        "  deck_title: \"Bài giảng\"\n"
        "  author: \"\"\n\n"
        "YÊU CẦU RÀNG BUỘC:\n"
        "- Tổng số slide: 5 đến 10.\n"
        "- Mỗi phần tử trong slides PHẢI có: slide (số thứ tự 1..n), title (3–8 từ), content (3–7 dòng văn bản, mỗi ý một dòng).\n"
        "- content sử dụng block scalar (|) với mỗi dòng cho một gạch đầu dòng; không dùng list lồng nhau.\n"
        "- Chỉ trả về YAML thuần, không code fence.\n"
        "- Ngôn ngữ: tiếng Việt.\n\n"
        "Nội dung cần chuyển:\n"
        f"{content_text}\n"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
        temperature=0.2,
    )
    yaml_text = resp.choices[0].message.content or ""
    # Loại bỏ code fences nếu có
    yaml_text = re.sub(r"^```[a-zA-Z]*\s*|\s*```$", "", yaml_text.strip(), flags=re.MULTILINE)

    def flatten_bullets(node):
        result = []
        if isinstance(node, list):
            for item in node:
                result.extend(flatten_bullets(item))
        elif isinstance(node, dict):
            for key, value in node.items():
                result.append(str(key).strip())
                result.extend(flatten_bullets(value))
        elif isinstance(node, str):
            result.append(re.sub(r"^\s*-\s+", "", node).strip())
        else:
            result.append(str(node).strip())
        return result

    try:
        data = yaml.safe_load(yaml_text) or {}
        slides = data.get("slides", [])
        normalized = []
        for idx, s in enumerate(slides, start=1):
            if "content" in s:
                content_text = str(s.get("content", "")).strip()
            else:
                bullets = flatten_bullets(s.get("bullets", []))
                content_text = "\n".join([line for line in bullets if line])
            normalized.append({
                "slide": s.get("slide", idx),
                "title": s.get("title", f"Slide {idx}"),
                "content": content_text
            })
        data["slides"] = normalized
        # Bảo toàn meta nếu có, nếu không thì thêm khung trống
        if "meta" not in data:
            data["meta"] = {"deck_title": "Bài giảng", "author": ""}
        yaml_text = yaml.safe_dump(data, allow_unicode=True, sort_keys=False)
    except Exception:
        # Nếu parse lỗi, vẫn lưu nguyên văn đã strip codefence
        pass

    # Tạo record mới trong collection content_yamls
    content_yaml_id = crepo.insert_content_yaml({
        "content_id": req.content_id,
        "yaml": yaml_text,
        "created_by": req.created_by,
    })
    return {"content_yaml_id": content_yaml_id, "yaml": yaml_text}


class ContentYAMLUpdateRequest(BaseModel):
    yaml: str
    updated_by: str | None = None


class ContentYAMLResponse(BaseModel):
    content_yaml_id: str
    content_id: str
    yaml: str
    created_at: str | None = None
    updated_at: str | None = None


@router.get("/template/yaml/{content_yaml_id}", response_model=ContentYAMLResponse)
def get_content_yaml(content_yaml_id: str):
    crepo = ContentRepository()
    doc = crepo.get_content_yaml_by_id(content_yaml_id)
    if not doc:
        raise HTTPException(status_code=404, detail="content_yaml_id not found")
    # Convert datetime to str
    for k in ("created_at", "updated_at"):
        if isinstance(doc.get(k), (str, type(None))):
            continue
        if doc.get(k):
            doc[k] = str(doc[k])
    return doc


@router.get("/template/yaml/by-content/{content_id}", response_model=list[ContentYAMLResponse])
def list_content_yaml_by_content(content_id: str):
    crepo = ContentRepository()
    docs = crepo.list_content_yaml_by_content(content_id)
    for doc in docs:
        for k in ("created_at", "updated_at"):
            if isinstance(doc.get(k), (str, type(None))):
                continue
            if doc.get(k):
                doc[k] = str(doc[k])
    return docs


@router.put("/template/yaml/{content_yaml_id}", response_model=ContentYAMLResponse)
def update_content_yaml(content_yaml_id: str, req: ContentYAMLUpdateRequest):
    crepo = ContentRepository()
    ok = crepo.update_content_yaml(content_yaml_id, req.yaml, updated_by=req.updated_by)
    if not ok:
        raise HTTPException(status_code=404, detail="content_yaml_id not found or not modified")
    # return updated doc
    return get_content_yaml(content_yaml_id)


@router.delete("/template/yaml/{content_yaml_id}")
def delete_content_yaml(content_yaml_id: str):
    crepo = ContentRepository()
    ok = crepo.delete_content_yaml(content_yaml_id)
    if not ok:
        raise HTTPException(status_code=404, detail="content_yaml_id not found")
    return {"status": "deleted", "content_yaml_id": content_yaml_id}


# ===================== Slide Templates (stored in DB via GridFS) =====================

class TemplateMetaResponse(BaseModel):
    template_id: str
    name: str
    filename: str
    content_type: str
    size: int | None = None
    description: str | None = None
    created_at: str | None = None
    updated_at: str | None = None


@router.post("/templates", response_model=TemplateMetaResponse)
async def upload_template(
    name: str = Form(...),
    description: str | None = Form(None),
    file: UploadFile = File(...),
):
    if not file.filename or not file.content_type:
        raise HTTPException(status_code=400, detail="Invalid file")
    data = await file.read()
    trepo = SlideTemplateRepository()
    trepo.create_indexes()
    template_id = trepo.insert_template(
        name=name,
        filename=file.filename,
        content_type=file.content_type,
        data=data,
        description=description,
    )
    doc = trepo.get_template_by_id(template_id)
    # Convert datetime to str
    for k in ("created_at", "updated_at"):
        if doc.get(k):
            doc[k] = str(doc[k])
    return doc


@router.get("/templates", response_model=list[TemplateMetaResponse])
def list_templates():
    trepo = SlideTemplateRepository()
    items = trepo.list_templates()
    for it in items:
        for k in ("created_at", "updated_at"):
            if it.get(k):
                it[k] = str(it[k])
    return items


@router.get("/templates/{template_id}", response_model=TemplateMetaResponse)
def get_template(template_id: str):
    trepo = SlideTemplateRepository()
    doc = trepo.get_template_by_id(template_id)
    if not doc:
        raise HTTPException(status_code=404, detail="template_id not found")
    for k in ("created_at", "updated_at"):
        if doc.get(k):
            doc[k] = str(doc[k])
    return doc


@router.get("/templates/{template_id}/download")
def download_template(template_id: str):
    trepo = SlideTemplateRepository()
    meta = trepo.get_template_by_id(template_id)
    if not meta:
        raise HTTPException(status_code=404, detail="template_id not found")
    data = trepo.download_template_file(template_id)
    if data is None:
        raise HTTPException(status_code=404, detail="template file not found")
    headers = {
        "Content-Disposition": f'attachment; filename="{meta.get("filename", "template.pptx")}"'
    }
    return Response(content=data, media_type=meta.get("content_type", "application/vnd.openxmlformats-officedocument.presentationml.presentation"), headers=headers)


@router.delete("/templates/{template_id}")
def delete_template(template_id: str):
    trepo = SlideTemplateRepository()
    ok = trepo.delete_template(template_id)
    if not ok:
        raise HTTPException(status_code=404, detail="template_id not found")
    return {"status": "deleted", "template_id": template_id}


# ===================== Template Inspection =====================

class PlaceholderInfo(BaseModel):
    name: str | None = None
    shape_type: str | None = None
    placeholder_type: str | None = None
    has_text: bool = False
    idx: int | None = None


class LayoutInfo(BaseModel):
    index: int
    name: str | None = None
    placeholders: list[PlaceholderInfo]


class TemplateInspectResponse(BaseModel):
    template_id: str
    filename: str | None = None
    layouts: list[LayoutInfo]


@router.get("/templates/{template_id}/inspect", response_model=TemplateInspectResponse)
def inspect_template(template_id: str):
    trepo = SlideTemplateRepository()
    meta = trepo.get_template_by_id(template_id)
    if not meta:
        raise HTTPException(status_code=404, detail="template_id not found")
    data = trepo.download_template_file(template_id)
    if not data:
        raise HTTPException(status_code=404, detail="template file not found")

    prs = Presentation(BytesIO(data))
    layouts: list[LayoutInfo] = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders: list[PlaceholderInfo] = []
        for shp in layout.shapes:
            info = PlaceholderInfo(
                name=getattr(shp, "name", None),
                shape_type=str(getattr(shp, "shape_type", None)),
                placeholder_type=None,
                has_text=getattr(shp, "has_text_frame", False),
                idx=getattr(shp, "placeholder_format", None).idx if hasattr(shp, "placeholder_format") else None,
            )
            if hasattr(shp, "placeholder_format"):
                try:
                    ph_type = shp.placeholder_format.type
                    info.placeholder_type = str(ph_type)
                except Exception:
                    pass
            placeholders.append(info)
        layouts.append(LayoutInfo(
            index=idx,
            name=getattr(layout, "name", None),
            placeholders=placeholders
        ))
    return TemplateInspectResponse(
        template_id=template_id,
        filename=meta.get("filename"),
        layouts=layouts,
    )
# ===================== Export PPTX from content_yaml_id + template_id =====================

class ExportPPTXRequest(BaseModel):
    content_yaml_id: str
    template_id: str
    filename: str | None = None
    overwrite_existing: bool | None = True


@router.post("/template/export")
def export_pptx(req: ExportPPTXRequest):
    """
    Tạo file PPTX từ content_yaml_id và template_id lưu trong DB.
    """
    # Load YAML
    crepo = ContentRepository()
    yaml_doc = crepo.get_content_yaml_by_id(req.content_yaml_id)
    if not yaml_doc:
        raise HTTPException(status_code=404, detail="content_yaml_id not found")
    yaml_text = yaml_doc.get("yaml", "")
    if not yaml_text:
        raise HTTPException(status_code=400, detail="YAML is empty for this content_yaml_id")

    def _try_load_yaml(text: str):
        try:
            return yaml.safe_load(text) or {}
        except Exception:
            return None

    data = _try_load_yaml(yaml_text)
    if data is None:
        # Fallback repair: quote all bullet items to avoid YAML nested block parsing issues
        lines = yaml_text.splitlines()
        repaired = []
        in_bullets = False
        bullets_indent = None
        for ln in lines:
            # Track entering bullets section
            m_bul = re.match(r'^(\s*)bullets\s*:\s*$', ln)
            if m_bul:
                in_bullets = True
                bullets_indent = len(m_bul.group(1))
                repaired.append(ln)
                continue
            # If we leave bullets section when indentation decreases and line not blank
            if in_bullets:
                if ln.strip() == "":
                    repaired.append(ln)
                    continue
                current_indent = len(ln) - len(ln.lstrip(' '))
                if current_indent <= bullets_indent and not ln.lstrip().startswith('-'):
                    in_bullets = False
                # For items under bullets:, ensure quoted strings
                if in_bullets and re.match(r'^\s*-\s+.*$', ln):
                    # Extract text after "- "
                    prefix, text = re.match(r'^(\s*-\s+)(.*)$', ln).groups()
                    text = text.strip()
                    # If already quoted, keep; else quote
                    if not (text.startswith('"') and text.endswith('"')) and not (text.startswith("'") and text.endswith("'")):
                        text = text.replace('"', '\\"')
                        ln = f'{prefix}"{text}"'
                    repaired.append(ln)
                    continue
            repaired.append(ln)
        yaml_text_repaired = "\n".join(repaired)
        data = _try_load_yaml(yaml_text_repaired)
        if data is None:
            # Second fallback: normalize nested bullet indentation inside bullets: block
            normalized = []
            in_bullets = False
            bullets_indent = None
            bullet_item_indent = None
            for ln in repaired:
                m_bul = re.match(r'^(\s*)bullets\s*:\s*$', ln)
                if m_bul:
                    in_bullets = True
                    bullets_indent = len(m_bul.group(1))
                    bullet_item_indent = None
                    normalized.append(ln)
                    continue
                if in_bullets:
                    if ln.strip() == "":
                        normalized.append(ln)
                        continue
                    current_indent = len(ln) - len(ln.lstrip(' '))
                    # Leaving bullets section
                    if current_indent <= bullets_indent and not ln.lstrip().startswith('-'):
                        in_bullets = False
                        bullets_indent = None
                        bullet_item_indent = None
                        normalized.append(ln)
                        continue
                    # Track first bullet item's indent
                    if re.match(r'^\s*-\s+', ln) and bullet_item_indent is None:
                        bullet_item_indent = len(re.match(r'^(\s*)-\s+', ln).group(1))
                        normalized.append(ln)
                        continue
                    # If a nested list item like:        - "  - text"
                    if bullet_item_indent is not None and re.match(r'^\s+-\s+"  - ', ln):
                        ln = ' ' * bullet_item_indent + ln.lstrip()
                        normalized.append(ln)
                        continue
                normalized.append(ln)
        yaml_text_normalized = "\n".join(normalized)
        data = _try_load_yaml(yaml_text_normalized)
        if data is None:
            def parse_relaxed(text: str) -> Dict[str, Any]:
                slides: List[Dict[str, Any]] = []
                meta: Dict[str, Any] = {}
                current: Optional[Dict[str, Any]] = None
                section = None
                bullet_mode = False
                bullet_indent = None

                def strip_quotes(val: str) -> str:
                    v = val.strip()
                    if len(v) >= 2 and ((v.startswith('"') and v.endswith('"')) or (v.startswith("'") and v.endswith("'"))):
                        return v[1:-1]
                    return v

                lines = text.splitlines()
                for raw in lines:
                    line = raw.rstrip('\r\n')
                    stripped = line.strip()
                    if not stripped:
                        continue
                    if stripped.startswith("#"):
                        continue

                    if re.match(r'^\s*slides\s*:\s*$', line):
                        section = "slides"
                        current = None
                        bullet_mode = False
                        bullet_indent = None
                        continue
                    if re.match(r'^\s*meta\s*:\s*$', line):
                        section = "meta"
                        current = None
                        bullet_mode = False
                        bullet_indent = None
                        continue

                    if section == "slides":
                        if bullet_mode:
                            m_bullet = re.match(r'^(\s*)-\s+(.*)$', line)
                            if m_bullet:
                                text_val = strip_quotes(m_bullet.group(2))
                                if current is not None:
                                    current.setdefault("bullets", []).append(text_val)
                                if bullet_indent is None:
                                    bullet_indent = len(m_bullet.group(1))
                                continue
                            else:
                                bullet_mode = False
                                bullet_indent = None
                                # fall through to process current line

                        m_layout = re.match(r'^\s*-\s*layout\s*:\s*(.+)$', line)
                        if m_layout:
                            current = {
                                "layout": strip_quotes(m_layout.group(1)),
                                "title": "",
                                "bullets": []
                            }
                            slides.append(current)
                            continue

                        if current is None:
                            continue

                        m_title = re.match(r'^\s*title\s*:\s*(.+)$', line)
                        if m_title:
                            current["title"] = strip_quotes(m_title.group(1))
                            continue

                        if re.match(r'^\s*bullets\s*:\s*$', line):
                            bullet_mode = True
                            bullet_indent = None
                            continue

                        # Any other property inside slide (ignored)
                        continue

                    if section == "meta":
                        m_kv = re.match(r'^\s*([A-Za-z0-9_]+)\s*:\s*(.+)$', line)
                        if m_kv:
                            meta[m_kv.group(1)] = strip_quotes(m_kv.group(2))
                        continue

                if not slides:
                    raise ValueError("No slides parsed in relaxed mode")
                if "deck_title" not in meta:
                    meta["deck_title"] = "Bài giảng"
                if "author" not in meta:
                    meta["author"] = ""
                return {"slides": slides, "meta": meta}

            try:
                data = parse_relaxed(yaml_text)
            except Exception:
                raise HTTPException(status_code=400, detail="Invalid YAML: could not parse after repair")

    # Load template PPTX bytes
    trepo = SlideTemplateRepository()
    meta = trepo.get_template_by_id(req.template_id)
    if not meta:
        raise HTTPException(status_code=404, detail="template_id not found")
    tpl_bytes = trepo.download_template_file(req.template_id)
    if not tpl_bytes:
        raise HTTPException(status_code=404, detail="template file not found")

    # Helper: flatten bullets support both flat strings and nested dict form -> returns list of plain strings
    def flatten_bullets(bullets):
        result = []
        for b in bullets or []:
            if isinstance(b, str):
                result.append(re.sub(r"^\s*-\s+", "", b).strip())
            elif isinstance(b, dict):
                for key, value in b.items():
                    result.append(str(key).strip())
                    result.extend(flatten_bullets(value))
            else:
                result.append(str(b).strip())
        return [line for line in result if line]

    def extract_lines(content: str) -> List[str]:
        if not content:
            return []
        lines = []
        for raw in content.splitlines():
            text = raw.strip()
            if not text:
                continue
            text = re.sub(r"^\s*-\s+", "", text)
            lines.append(text)
        return lines

    # Build PPTX in-memory
    prs = Presentation(BytesIO(tpl_bytes))

    def _find_title_body_placeholders(slide):
        """
        Ưu tiên theo nội dung mặc định của Canva:
        - "Add a heading"  -> title
        - "Add a little bit of body text" -> body
        Sau đó theo idx (0 = title, 1 = body), rồi theo type, cuối cùng fallback quét shapes.
        Hỗ trợ trường hợp content là OBJECT (7) thay vì BODY (2).
        """
        title_tf = None
        body_tf = None

        # 0) Ưu tiên dò theo text mặc định
        for shp in slide.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            tf = shp.text_frame
            if not tf:
                continue
            try:
                current_text = (tf.text or "").strip()
            except Exception:
                current_text = ""
            if current_text == "Add a heading" and title_tf is None:
                title_tf = tf
            elif current_text == "Add a little bit of body text" and body_tf is None:
                body_tf = tf
            if title_tf is not None and body_tf is not None:
                return title_tf, body_tf

        # 1) Ưu tiên theo idx
        for ph in getattr(slide.shapes, "placeholders", []):
            if not getattr(ph, "has_text_frame", False):
                continue
            try:
                idx = ph.placeholder_format.idx
            except Exception:
                idx = None
            if idx == 0 and title_tf is None:
                title_tf = ph.text_frame
            elif idx == 1 and body_tf is None:
                body_tf = ph.text_frame
        if title_tf is not None and body_tf is not None:
            return title_tf, body_tf

        # 2) Theo type (bao quát BODY/OBJECT/SUBTITLE)
        for ph in getattr(slide.shapes, "placeholders", []):
            if not getattr(ph, "has_text_frame", False):
                continue
            try:
                ph_type = ph.placeholder_format.type
            except Exception:
                continue
            if title_tf is None and ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                title_tf = ph.text_frame
            if body_tf is None:
                if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE):
                    body_tf = ph.text_frame
                elif hasattr(PP_PLACEHOLDER, "OBJECT") and ph_type == getattr(PP_PLACEHOLDER, "OBJECT"):
                    body_tf = ph.text_frame
            if title_tf is not None and body_tf is not None:
                break

        # 3) Fallback: shapes.title và shape đầu tiên có text (không phải title)
        if title_tf is None:
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape and getattr(title_shape, "has_text_frame", False):
                title_tf = title_shape.text_frame
        if body_tf is None:
            for shp in slide.shapes:
                if not getattr(shp, "has_text_frame", False):
                    continue
                if shp == getattr(slide.shapes, "title", None):
                    continue
                body_tf = shp.text_frame
                break

        return title_tf, body_tf

    # Optional cover slide
    meta_yaml = data.get("meta", {})
    deck_title = meta_yaml.get("deck_title")
    author = meta_yaml.get("author")
    num_existing = len(prs.slides)
    if deck_title or author:
        try:
            if req.overwrite_existing and num_existing >= 1:
                cover = prs.slides[0]
            else:
                cover = prs.slides.add_slide(prs.slide_layouts[0])
            title_tf, body_tf = _find_title_body_placeholders(cover)
            if deck_title and title_tf is not None:
                title_tf.clear()
                p = title_tf.paragraphs[0]
                p.text = deck_title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.name = "Arial"
            if author and body_tf is not None:
                body_tf.clear()
                p = body_tf.paragraphs[0]
                p.text = author
                p.font.size = Pt(24)
                p.font.name = "Arial"
        except Exception:
            # If template lacks title slide, skip
            pass

    for i, s in enumerate(data.get("slides", [])):
        title = s.get("title", "")
        if "content" in s:
            bullet_lines = extract_lines(s.get("content", ""))
        else:
            bullet_lines = flatten_bullets(s.get("bullets", []))
        # Luôn dùng layout "Title and Content" (index 1).
        # Nếu overwrite_existing=True và đã có sẵn slide ở vị trí i+1, ghi đè vào đó.
        if req.overwrite_existing and len(prs.slides) > (1 + i):
            slide = prs.slides[1 + i]
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Ưu tiên chọn placeholder theo idx chuẩn của layout này: title=0, body=1
        title_tf = None
        body_tf = None
        try:
            for ph in slide.shapes.placeholders:
                try:
                    idx = ph.placeholder_format.idx
                except Exception:
                    idx = None
                if idx == 0 and getattr(ph, "has_text_frame", False):
                    title_tf = ph.text_frame
                elif idx == 1 and getattr(ph, "has_text_frame", False):
                    body_tf = ph.text_frame
        except Exception:
            pass
        # Nếu chưa thấy, fallback dò theo loại/shape
        if title_tf is None or body_tf is None:
            t2, b2 = _find_title_body_placeholders(slide)
            if title_tf is None:
                title_tf = t2
            if body_tf is None:
                body_tf = b2

        if title_tf is not None:
            title_tf.clear()
            p = title_tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = "Arial"
        if body_tf is not None:
            body_tf.clear()
            first = True
            for text in bullet_lines:
                if first:
                    p = body_tf.paragraphs[0]
                    first = False
                else:
                    p = body_tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(20)
                p.font.name = "Arial"

    out = BytesIO()
    prs.save(out)
    out.seek(0)

    # Tên file: đảm bảo có đuôi .pptx và header theo RFC 5987 để hỗ trợ Unicode
    filename = req.filename or f"slides_{req.content_yaml_id}.pptx"
    if not filename.lower().endswith(".pptx"):
        filename = f"{filename}.pptx"
    from urllib.parse import quote
    ascii_fallback = "slides.pptx"
    try:
        filename.encode("latin-1")
        content_disp = f'attachment; filename="{filename}"'
    except UnicodeEncodeError:
        # Dùng filename* cho UTF-8 theo RFC 5987 + fallback ascii
        content_disp = f'attachment; filename="{ascii_fallback}"; filename*=UTF-8\'\'{quote(filename)}'
    headers = {"Content-Disposition": content_disp}
    return Response(
        content=out.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )

